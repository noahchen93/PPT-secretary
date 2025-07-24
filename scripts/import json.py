import json
import os
import time
import requests
from pptx import Presentation
from pptx.util import Inches
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.options import Options

# --------------------------------------------------------------------------
# 你从Pexels获取的API Key
# --------------------------------------------------------------------------
PEXELS_API_KEY = "WYTIJDKwaIxU8DujqrYnnTb8owdBSEXTiiXp1luyQgUBa3JC2rw8CU1C"
# --------------------------------------------------------------------------

# 为幻灯片设计的HTML和CSS模板
HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <title>Slide</title>
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@400;700&display=swap');
        body {{ margin: 0; font-family: 'Noto Sans SC', sans-serif; }}
        .slide {{
            width: 1280px;
            height: 720px;
            display: flex;
            background-color: #ffffff;
            overflow: hidden;
        }}
        .content-pane {{
            width: 58%;
            padding: 50px 40px 50px 80px;
            box-sizing: border-box;
            display: flex;
            flex-direction: column;
            justify-content: center;
        }}
        .title {{
            font-size: 44px;
            font-weight: 700;
            color: #0d47a1; /* 深蓝色 */
            margin-bottom: 35px;
            line-height: 1.3;
        }}
        .points ul {{
            list-style: none;
            padding-left: 0;
            margin: 0;
        }}
        .points li {{
            font-size: 24px;
            color: #37474f; /* 深灰色 */
            margin-bottom: 18px;
            padding-left: 35px;
            position: relative;
            line-height: 1.6;
        }}
        .points li::before {{
            content: '◆';
            position: absolute;
            left: 0;
            color: #1976d2; /* 亮蓝色 */
            font-size: 20px;
            top: 5px;
        }}
        .image-pane {{
            width: 42%;
            background-size: cover;
            background-position: center;
            /* 倾斜的剪裁路径 */
            clip-path: polygon(25% 0, 100% 0, 100% 100%, 0% 100%);
        }}
    </style>
</head>
<body>
    <div class="slide" id="capture">
        <div class="content-pane">
            <div class="title">{title}</div>
            <div class="points"><ul>{points_html}</ul></div>
        </div>
        <div class="image-pane" style="background-image: url('{image_url}');"></div>
    </div>
</body>
</html>
"""

def get_image_url(query, api_key):
    """使用Pexels API搜索图片并返回图片的URL。"""
    if not api_key or api_key == "YOUR_API_KEY_HERE":
        return "https://images.pexels.com/photos/356079/pexels-photo-356079.jpeg" # 默认图片

    headers = {"Authorization": api_key}
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1&orientation=landscape"
    try:
        response = requests.get(url, headers=headers, timeout=15)
        response.raise_for_status()
        data = response.json()
        if data.get("photos"):
            return data["photos"][0]["src"]["large2x"]
    except requests.exceptions.RequestException as e:
        print(f"网络错误: {e}")
    return "https://images.pexels.com/photos/356079/pexels-photo-356079.jpeg" # 默认图片

def create_presentation_from_web(json_data_str, output_filename):
    """通过生成网页截图来创建高质量排版的PPT。"""
    # 将所有文件生成到你指定的文件夹路径
    output_dir = r"C:\Users\Noah Chen\Documents\BaiduSyncdisk\个人文档\VS Creation"
    os.makedirs(output_dir, exist_ok=True)
    print(f"所有文件将被保存在 '{output_dir}/' 文件夹中。")

    print("正在初始化虚拟浏览器...")
    chrome_options = Options()
    chrome_options.add_argument("--headless") # 无头模式，不在屏幕上显示浏览器
    chrome_options.add_argument("--hide-scrollbars")
    chrome_options.add_argument("--disable-gpu")
    chrome_options.add_argument("--window-size=1280,720")
    
    # 使用webdriver-manager自动安装和管理chromedriver
    service = Service(ChromeDriverManager().install())
    driver = webdriver.Chrome(service=service, options=chrome_options)
    
    data = json.loads(json_data_str)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # --- 标题幻灯片 (传统方式创建) ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.get("title", "")
    slide.placeholders[1].text = data.get("subtitle", "")

    temp_files = []

    # --- 内容幻灯片 (通过网页截图创建) ---
    for i, item in enumerate(data.get("slides", [])):
        print(f"正在处理第 {i+1}/{len(data.get('slides',[]))}: {item['title']}")
        
        # 1. 获取图片URL
        image_url = get_image_url(item.get("search_term"), PEXELS_API_KEY)
        
        # 2. 生成HTML内容
        points_html = "".join([f"<li>{p}</li>" for p in item.get("points", [])])
        slide_html = HTML_TEMPLATE.format(title=item.get("title"), points_html=points_html, image_url=image_url)
        
        # 3. 创建临时HTML文件
        html_filename = os.path.join(output_dir, f"temp_slide_{i}.html")
        with open(html_filename, "w", encoding="utf-8") as f:
            f.write(slide_html)
        temp_files.append(html_filename)
        
        # 4. 使用Selenium截图
        driver.get("file:///" + os.path.abspath(html_filename))
        time.sleep(1.5) # 增加一点等待时间确保图片完全渲染
        screenshot_filename = os.path.join(output_dir, f"screenshot_{i}.png")
        driver.find_element(By.ID, "capture").screenshot(screenshot_filename)
        temp_files.append(screenshot_filename)
        
        # 5. 将截图添加到PPT
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Layout 6 是空白布局
        slide.shapes.add_picture(screenshot_filename, Inches(0), Inches(0), width=prs.slide_width)

    # --- 总结幻灯片 (传统方式创建) ---
    summary_data = data.get("summary")
    if summary_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = summary_data.get("title", "总结")
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for point in summary_data.get("points", []):
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

    driver.quit()
    final_ppt_path = os.path.join(output_dir, output_filename)
    print(f"\n正在保存PPT到 {final_ppt_path}...")
    prs.save(final_ppt_path)
    
    # --- 清理临时文件 ---
    print("正在清理临时文件...")
    for f in temp_files:
        os.remove(f)
        
    print(f"\n🎉 恭喜！排版精美的演示文稿 '{final_ppt_path}' 已成功创建！")

if __name__ == '__main__':
    json_content = """
    {
      "title": "人工智能的未来：机遇与挑战",
      "subtitle": "探索下一代AI技术",
      "slides": [
        { "title": "1. AI大模型的演进", "search_term": "abstract technology background", "points": ["从单一功能到通用智能（AGI）的潜力", "多模态能力：理解文本、图像、声音和视频", "模型规模与效率的平衡"] },
        { "title": "2. AI在各行业的应用前景", "search_term": "futuristic medical technology", "points": ["医疗健康：个性化诊断与药物研发", "自动驾驶：提升交通安全与效率", "创意产业：自动化内容生成与辅助设计"] },
        { "title": "3. 面临的挑战与道德考量", "search_term": "data privacy ethics", "points": ["数据隐私与安全问题", "算法的公平性与偏见", "对就业市场和社会结构的冲击"] }
      ],
      "summary": { "title": "总结与展望", "points": ["AI是推动社会变革的关键力量", "技术发展需与伦理规范并行", "拥抱合作，共创负责任的AI未来"] }
    }
    """
    create_presentation_from_web(json_content, "人工智能的未来_网页版.pptx")
